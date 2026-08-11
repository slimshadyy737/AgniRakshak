"""
AgniRakshak PowerPoint Generator (.pptx)
Generates official 9-slide PowerPoint deck for IIC 3.0 Open Innovation track.
"""

import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # Set slide dimensions to widescreen 16:9 (13.33 x 7.5 inches)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]
    
    # Theme Colors
    COLOR_BG = RGBColor(15, 23, 42)        # Slate Dark 900
    COLOR_CARD = RGBColor(30, 41, 59)      # Slate 800
    COLOR_ORANGE = RGBColor(249, 115, 22)  # Orange 500
    COLOR_WHITE = RGBColor(248, 250, 252)  # Slate 50
    COLOR_MUTED = RGBColor(148, 163, 184) # Slate 400
    COLOR_GREEN = RGBColor(16, 185, 129)  # Emerald 500

    slides_data = [
        {
            "slide_num": "Slide 1",
            "title": "IDEA TITLE & TEAM DETAILS",
            "subtitle": "AgniRakshak: Distributed Edge-AI Wildfire Early Detection Network",
            "content": [
                ("Theme / Track", "Open Innovation"),
                ("Idea Title", "AgniRakshak (अग्निरक्षक)"),
                ("Tagline", "Detect Early. Alert Faster. Prevent Wildfires."),
                ("Team Name", "Hell Fire Club"),
                ("Institution", "Manipal University Jaipur (IIC 3.0)"),
                ("Core Value Prop", "Sub-minute wildfire ignition detection combining solar-powered IoT micro-climate sensor fusion with temporal derivative AI classification (dT/dt, dCO/dt).")
            ]
        },
        {
            "slide_num": "Slide 2",
            "title": "PROBLEM STATEMENT",
            "subtitle": "Critical Latency & Failure of Existing Monitoring Methods",
            "content": [
                ("Satellite & Camera Latency", "Satellites have 1-3 hour pass delays; by detection time, small ignitions escalate into catastrophic mega-fires."),
                ("Remote Forest Blindspots", "Deep wildlands lack power grid infrastructure and continuous physical human surveillance."),
                ("Single-Threshold Failure", "Conventional static temperature alarms produce false triggers from hot ambient weather or fail in outdoor wind."),
                ("Severe Consequences", "Massive ecological destruction, megatons of unmitigated CO2 emissions, and loss of life & property.")
            ]
        },
        {
            "slide_num": "Slide 3",
            "title": "PROPOSED SOLUTION",
            "subtitle": "Multi-Sensor Edge-AI Sensor Fusion Platform",
            "content": [
                ("Distributed Sensor Nodes", "Deploy autonomous micro-climate nodes capturing Temp, Humidity, CO, Smoke, and Barometric Pressure."),
                ("Temporal Derivative Engine", "Computes rate-of-change metrics (dT/dt, dCO/dt) to detect chemical ignition before visible smoke/flames."),
                ("Hybrid AI Classifier", "Scikit-Learn Random Forest model + deterministic heuristic safety rules classify risk: NORMAL, WARNING, HIGH RISK."),
                ("GPS & Sub-Minute Alerts", "Instant geolocation tagging via LoRa / Wi-Fi mesh dispatches responders within 60 seconds.")
            ]
        },
        {
            "slide_num": "Slide 4",
            "title": "INNOVATION & UNIQUENESS",
            "subtitle": "Why AgniRakshak Outperforms Traditional Fire Alarms",
            "content": [
                ("Feature Fusion vs Static Thresholds", "Combines multi-parameter temporal trends instead of single threshold triggers (e.g. Temp > X)."),
                ("Sub-Minute Detection Speed", "Local edge detection triggers instant alarm buzzer and cloud dashboard alerts."),
                ("Physics-Based Calibration", "Polynomial bivariate temperature & humidity correction K(T, H) eliminates sensor drift false alarms."),
                ("Ultra-Low Cost & Off-Grid", "Complete node BOM < $35 with solar MPPT harvesting and 19-day dark autonomy.")
            ]
        },
        {
            "slide_num": "Slide 5",
            "title": "TECHNICAL APPROACH & ARCHITECTURE",
            "subtitle": "Hardware, Firmware, and Cloud AI Stack",
            "content": [
                ("Hardware Layer", "ESP32 DevKit MCU + Bosch BME280 + MQ-7 CO (5V/1.4V cycle) + MQ-2 Smoke + NEO-6M GPS."),
                ("Firmware Architecture", "FreeRTOS dual-core task scheduling (Core 0: Sensor conditioning; Core 1: OLED, Alerts, Wi-Fi TX)."),
                ("AI & Software Stack", "Python Streamlit Web Dashboard + Scikit-Learn Random Forest Model + Folium Geospatial Map."),
                ("Power Management", "5W Monocrystalline Solar + CN3791 MPPT + 6000mAh LiFePO4 battery (< 13mA avg draw).")
            ]
        },
        {
            "slide_num": "Slide 6",
            "title": "PROTOTYPE & WORKING DEMONSTRATION",
            "subtitle": "Full-Stack Implementation Deliverables",
            "content": [
                ("Interactive Streamlit Dashboard", "Real-time web monitoring interface featuring live risk gauges, Folium map, and Plotly charts."),
                ("Live Scenario Simulator", "Test controls for 4 scenarios: Normal Ambient, Hot & Dry Weather, Smoldering Fire, Active Wildfire."),
                ("ESP32 Firmware Code", "Production C++/Arduino sketch with FreeRTOS, OLED status display, active buzzer, and JSON HTTP client."),
                ("Trained ML Engine", "Persisted Scikit-Learn Random Forest classifier model artifacts (wildfire_model.pkl).")
            ]
        },
        {
            "slide_num": "Slide 7",
            "title": "FEASIBILITY & VIABILITY",
            "subtitle": "Economic & Power System Autonomy",
            "content": [
                ("Node Cost (< $35)", "ESP32 ($4.50), BME280 ($3.20), MQ-7/MQ-2 ($4.80), GPS & OLED ($6.50), Solar MPPT & Battery ($12.00)."),
                ("5-Min Periodic Duty Cycle", "Consumes 12.63 mA average current at 3.3V (0.041 W power consumption)."),
                ("19-Day Dark Autonomy", "6000 mAh LiFePO4 battery runs over 19 days without solar replenishment during dense smoke plumes."),
                ("IP65 Polycarbonate Enclosure", "Dual-compartment sealed electronics housing with louvered 45° rain protection sensing chamber.")
            ]
        },
        {
            "slide_num": "Slide 8",
            "title": "IMPACT & FUTURE ROADMAP",
            "subtitle": "Environmental Preservation & Scalability",
            "content": [
                ("Ecological Impact", "Protects forest biodiversity, wildlife habitats, and reduces unmitigated carbon emissions."),
                ("Economic Protection", "Prevents multi-million dollar timber, agricultural, and community infrastructure losses."),
                ("LoRa Mesh Scaling", "Expand nodes across thousands of forest hectares using sub-GHz LoRa RF mesh networking."),
                ("Future Integration", "Combine ground IoT sensor mesh telemetry with autonomous drone verification and satellite maps.")
            ]
        },
        {
            "slide_num": "Slide 9",
            "title": "CONCLUSION & SUMMARY VISION",
            "subtitle": "Detect Early. Alert Faster. Prevent Wildfires.",
            "content": [
                ("Core Vision", "Deploy a continuous, autonomous network of edge-AI nodes across vulnerable forest ecosystems."),
                ("System Workflow", "MONITOR -> DETECT -> ANALYZE -> ALERT -> ACT"),
                ("Competitive Advantage", "Open Innovation track choice preserves Round 1 solution for offline physical presentation at MUJ."),
                ("Final Tagline", "AI + IoT + Early Sensor Fusion = A Smarter Approach to Wildfire Prevention.")
            ]
        }
    ]

    for slide_info in slides_data:
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add Full Dark Background Shape
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background()
        
        # Add Header Title Box
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.2))
        tf = title_box.text_frame
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = f"{slide_info['slide_num']} — {slide_info['title']}"
        p0.font.bold = True
        p0.font.size = Pt(26)
        p0.font.color.rgb = COLOR_ORANGE
        
        p1 = tf.add_paragraph()
        p1.text = slide_info['subtitle']
        p1.font.size = Pt(16)
        p1.font.color.rgb = COLOR_MUTED
        
        # Add Content Cards
        left_margin = Inches(0.8)
        top_start = Inches(1.8)
        card_width = Inches(11.7)
        card_height = Inches(5.0)
        
        content_box = slide.shapes.add_textbox(left_margin, top_start, card_width, card_height)
        ctf = content_box.text_frame
        ctf.word_wrap = True
        
        for idx, (label, desc) in enumerate(slide_info['content']):
            p = ctf.add_paragraph() if idx > 0 else ctf.paragraphs[0]
            p.text = f"• {label}: "
            p.font.bold = True
            p.font.size = Pt(16)
            p.font.color.rgb = COLOR_GREEN
            
            # Add description text
            run = p.add_run()
            run.text = desc
            run.font.bold = False
            run.font.size = Pt(15)
            run.font.color.rgb = COLOR_WHITE
            p.space_after = Pt(14)
            
    output_path = "AgniRakshak_IIC3_Presentation.pptx"
    prs.save(output_path)
    print(f"PowerPoint Presentation successfully created at {output_path}")

if __name__ == "__main__":
    create_presentation()
